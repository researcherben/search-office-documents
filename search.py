#!/usr/bin/env python3

import glob # https://docs.python.org/3/library/glob.html
import argparse  # https://docs.python.org/3.3/library/argparse.html
import docx
import pptx
import re # https://docs.python.org/3/library/re.html

"""
https://automatetheboringstuff.com/chapter13/

options for searching document content:
* use pandoc to convert .docx/.pptx/.pdf to plain text files, then search text files. 
    (-) This will result in loss of "which page is the content on"

* use python-docx 
    (-) Not all text in the document appears in "paragraph" so extra complexity needed to get all text
    (-) doesn't support page numbers
https://python-docx.readthedocs.io/en/latest/user/documents.html

https://python-pptx.readthedocs.io/en/latest/

LATER
extract headers too: https://stackoverflow.com/a/54281094/1164295
https://python-docx.readthedocs.io/en/latest/user/hdrftr.html
"""

def open_docx_file(path_to_docx: str) -> str:
    """
    https://stackoverflow.com/a/35871416/1164295

    https://python-docx.readthedocs.io/en/latest/user/documents.html
    """
    doc = docx.Document(path_to_docx)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def search_docx(this_regex_search_string:str):
    """
    """
    list_of_files = glob.glob("**/*.docx", recursive=True)

    match_count = 0
    for this_filename in list_of_files:
        file_content = open_docx_file(this_filename)
        result = re.search(this_regex_search_string, file_content)
        if result: # no match results in None
            print(this_filename, "has the string", result.group(0))
            match_count+=1
    print("\nthere were",match_count,"files that contain",args.regex_search_string)


def search_pptx(this_regex_search_string:str):
    """
    """
    list_of_files = glob.glob("**/*.pptx", recursive=True)

    match_count = 0
    for this_filename in list_of_files:
        file_content = open_pptx_file(this_filename)
        result = re.search(this_regex_search_string, file_content)
        if result: # no match results in None
            print(this_filename, "has the string", result.group(0))
            match_count+=1
    print("\nthere were",match_count,"files that contain",args.regex_search_string)


def open_pptx_file(path_to_pptx: str) -> str:
    """
    https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
    """
    prs = pptx.Presentation(path_to_pptx)

    text_runs = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs += run.text
    return text_runs


def arg_parse():
    theparser = argparse.ArgumentParser(
        description="search a directory for a regex", allow_abbrev=False
    )

    # ********** begin argparse configuration *****************
    theparser.add_argument(
        "folder_path",
        metavar="folder_path", # this is what appears in the -h output
        type=str,
        default=".",
        help="Directory to search. Required. Default is '.'",
    )
    theparser.add_argument(
        "regex_search_string",
        metavar="regex_search_string", # this is what appears in the -h output
        type=str,
        default="dog",
        help="regex for string search. Default is 'dog'",
    )

    # ********** end argparse configuration *****************
    return theparser


if __name__ == "__main__":

    theparser = arg_parse()
    args = theparser.parse_args()
    #print("args=", args.folder_path, args.regex_search_string)

    search_docx(args.regex_search_string)

    list_of_files = glob.glob("**/*.pptx", recursive=True)
    for this_filename in list_of_files:
        out = open_pptx_file(this_filename)
        print(out)


# EOF
